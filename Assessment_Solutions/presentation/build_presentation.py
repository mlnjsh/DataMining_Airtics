"""
Builds Formative_Assessment_Presentation.pptx
13 content slides + cover + references + thank-you, with speaker notes.
Run:  python build_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

IMG = "img"
OUT = "Formative_Assessment_Presentation.pptx"

# ---------- palette ----------
NAVY   = RGBColor(0x1F, 0x3A, 0x5F)
BLUE   = RGBColor(0x2E, 0x5A, 0x88)
RED    = RGBColor(0xD0, 0x34, 0x2C)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x22, 0x22, 0x22)
FONT   = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height

slide_no = [0]


# ---------- helpers ----------
def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items, size=18):
    """items: list of str or (str, level) tuples; '**' wraps bold."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, level = (item, 0) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(6)
        bullet = "• " if level == 0 else "– "
        parts = text.split("**")
        first = True
        for j, part in enumerate(parts):
            if part == "":
                continue
            r = p.add_run()
            r.text = (bullet if first else "") + part
            first = False
            r.font.size = Pt(size - 2 * level)
            r.font.bold = (j % 2 == 1)
            r.font.color.rgb = BLACK
            r.font.name = FONT
    return tb


def add_code(slide, left, top, width, height, code, size=13):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = RGBColor(0xC9, 0xD3, 0xDF)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.margin_top = tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.name = "Consolas"
        r.font.color.rgb = NAVY if line.strip().startswith("#") else BLACK
    return box


def add_table(slide, left, top, width, height, rows, col_widths=None, size=13):
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = str(val)
            run.font.size = Pt(size); run.font.name = FONT
            run.font.bold = (r == 0)
            run.font.color.rgb = WHITE if r == 0 else BLACK
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r == 0 else (LIGHT if r % 2 == 0 else WHITE)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
    return shape


def add_image(slide, path, left, top, width=None, height=None):
    return slide.shapes.add_picture(os.path.join(IMG, path), left, top, width=width, height=height)


def content_slide(title, section=None, notes=""):
    slide = prs.slides.add_slide(BLANK)
    slide_no[0] += 1
    # title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.15), Inches(11), Inches(0.8), title,
             size=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if section:
        add_text(slide, Inches(9.3), Inches(0.15), Inches(3.6), Inches(0.8), section,
                 size=13, color=RGBColor(0xC9, 0xD3, 0xDF), align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)
    # footer
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.35),
             "Linear Regression Analysis  |  MSc Data Science  |  Formative Assignment",
             size=10, color=GRAY)
    add_text(slide, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.35),
             str(slide_no[0]), size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


# =====================================================================
#  COVER
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.6), W, Inches(0.08))
strip.fill.solid(); strip.fill.fore_color.rgb = RED; strip.line.fill.background()
add_text(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.6),
         "Linear Regression Analysis", size=48, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.9), Inches(11.5), Inches(1.4),
         "Concepts, Types, and Comparison with Logistic Regression",
         size=26, color=RGBColor(0xC9, 0xD3, 0xDF))
add_text(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(2),
         ["Master of Science in Data Science  |  Formative Assignment (40 marks)",
          "Student Name:  ____________________        Student ID:  ____________",
          "Email:  ____________________               Date:  ____________"],
         size=16, color=WHITE)
s.notes_slide.notes_text_frame.text = (
    "Good morning. This presentation covers three things: what linear regression is, "
    "the different types of linear regression with examples from different professions, "
    "and how linear regression differs from logistic regression. All R output shown was "
    "produced by running the code on the data sets named on each slide.")

# =====================================================================
#  PART 1 : CONCEPT  (slides 1-4)
# =====================================================================
s = content_slide("What is Linear Regression?", "Part 1 : Concept",
    "Linear regression answers two questions. First, how strongly is y related to x, and in "
    "which direction? Second, what value of y should we expect for a new x? The word linear "
    "means the relationship is modelled as a straight line. b0 is where the line crosses the "
    "y axis and b1 is how steep it is. The error term epsilon is everything the line cannot "
    "explain. Regression is used for two purposes: explanation (which factors matter and by how "
    "much) and prediction (what will y be for a new case).")
add_bullets(s, Inches(0.6), Inches(1.4), Inches(6.6), Inches(5.3), [
    "A statistical method that models the relationship between a **numeric outcome** (y) and one or more **predictors** (x)",
    "Fits a straight line through the data:",
    ("y  =  b₀  +  b₁ x  +  ε", 1),
    ("b₀ = intercept : value of y when x = 0", 1),
    ("b₁ = slope : change in y for a one-unit increase in x", 1),
    ("ε = error : what the line cannot explain", 1),
    "Two uses: **explain** which factors matter, and **predict** y for new cases",
    "Outcome must be **continuous** (price, weight, yield, mpg)",
], size=18)
add_image(s, "least_squares.png", Inches(7.4), Inches(1.5), width=Inches(5.5))

s = content_slide("How the Line is Fitted and Judged", "Part 1 : Concept",
    "The line is chosen by the method of least squares: it minimises the sum of the squared "
    "vertical distances between the points and the line. Those distances are the residuals. "
    "R-squared of 0.75 means 75 percent of the ups and downs in y are explained by x. The "
    "p-value of the slope tests the null hypothesis that the slope is zero; a p-value below "
    "0.05 means x genuinely helps explain y. Before trusting a model we check the four LINE "
    "assumptions using a residual plot and a Q-Q plot. If they fail, the p-values and "
    "predictions cannot be trusted.")
add_bullets(s, Inches(0.6), Inches(1.4), Inches(6.2), Inches(5.3), [
    "**Least squares**: choose b₀, b₁ to minimise  Σ (yᵢ − ŷᵢ)²",
    "Key outputs to read:",
    ("**Coefficients** with p-values : is the slope different from 0 ?", 1),
    ("**R²** : share of variation in y explained (0 to 1)", 1),
    ("**Residual standard error** : typical size of a prediction error", 1),
    "Assumptions (**LINE**):",
    ("**L**inearity  •  **I**ndependence", 1),
    ("**N**ormal residuals  •  **E**qual variance", 1),
    "Check with residual plots and a Q-Q plot before trusting the model",
], size=18)
add_table(s, Inches(7.2), Inches(1.6), Inches(5.6), Inches(2.4), [
    ["Output", "Meaning", "Good sign"],
    ["Slope p-value", "H₀: slope = 0", "< 0.05"],
    ["R²", "Variation explained", "Close to 1"],
    ["Residual SE", "Typical error", "Small vs range of y"],
    ["F-statistic", "Whole model useful?", "p < 0.05"],
], col_widths=[1.5, 2.2, 1.9], size=13)
add_code(s, Inches(7.2), Inches(4.3), Inches(5.6), Inches(1.6),
         "model <- lm(y ~ x, data = mydata)\nsummary(model)      # coefficients, R2, p-values\nplot(model)         # 4 diagnostic plots\npredict(model, newdata = data.frame(x = 5))")

s = content_slide("Example 1 : Car Weight and Fuel Economy", "Part 1 : Concept",
    "This example uses the mtcars data set that ships with R: 32 cars from a 1974 magazine. "
    "The relationship is negative, so the slope is below zero. Every extra thousand pounds "
    "costs about 5.3 miles per gallon, and weight alone explains 75 percent of the variation "
    "in fuel economy. The intercept of 37 is the predicted mpg of a car weighing nothing, which "
    "is physically meaningless but needed to position the line. Prediction example: a 3000 "
    "pound car, so wt equals 3, is predicted to do 37.3 minus 5.3 times 3, about 21 mpg. The "
    "p-value of the slope is 1.3 times 10 to the minus 10, so the effect is highly significant.")
add_image(s, "mpg_wt.png", Inches(0.5), Inches(1.4), width=Inches(6.6))
add_code(s, Inches(7.4), Inches(1.4), Inches(5.5), Inches(2.2),
         "model <- lm(mpg ~ wt, data = mtcars)\ncoef(model)\n# (Intercept)          wt\n#      37.285      -5.344\nsummary(model)$r.squared\n# 0.753")
add_bullets(s, Inches(7.4), Inches(4.1), Inches(5.5), Inches(2.8), [
    "Engineering data: 32 cars, weight vs miles per gallon",
    "Every extra 1000 lbs costs about **5.3 mpg**",
    "Weight explains **75 %** of fuel economy (R² = 0.75)",
    "Slope p-value = 1.3 × 10⁻¹⁰ : highly significant",
    "Prediction: a 3000 lb car → 37.3 − 5.3 × 3 ≈ **21 mpg**",
], size=16)

s = content_slide("More Examples from Different Fields", "Part 1 : Concept",
    "The same equation works in every field. What changes is the meaning of the slope: rupees "
    "of spending per rupee of income, millimetres of mercury per year of age, marks per hour of "
    "study, price per square foot. In economics the slope of spending on income is the marginal "
    "propensity to consume, typically 0.7 to 0.9. In medicine systolic blood pressure rises "
    "roughly half to one millimetre of mercury per year of adult age. Ask the audience: what is "
    "y and what is x in your own job? That is the best check that the idea has landed.")
add_table(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(4.2), [
    ["Field", "Outcome  (y)", "Predictor  (x)", "What the slope means"],
    ["Economics", "Household spending", "Household income", "Marginal propensity to consume, about 0.7 to 0.9"],
    ["Health", "Systolic blood pressure", "Age", "About 0.5 to 1 mmHg increase per year of adult age"],
    ["Education", "Final exam score", "Hours studied", "Extra marks per hour of study"],
    ["Real estate", "House price", "Floor area (sq ft)", "Price per square foot"],
    ["Agriculture", "Crop yield", "Fertilizer (kg / ha)", "Extra bushels per kilogram of fertilizer"],
    ["Marketing", "Monthly sales", "Advertising spend", "Sales generated per rupee of advertising"],
], col_widths=[1.7, 2.9, 2.7, 4.8], size=14)
add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.8),
         "Same equation every time  —  only the meaning and units of the slope change.",
         size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# =====================================================================
#  PART 2 : TYPES  (slides 5-9)
# =====================================================================
s = content_slide("Types of Linear Regression : Overview", "Part 2 : Types",
    "All of these are called linear because the outcome is a linear combination of the "
    "coefficients. Polynomial regression fits a curve but is solved by exactly the same least "
    "squares machinery. Ridge and Lasso are multiple regression with a penalty that shrinks the "
    "coefficients to prevent overfitting when there are many predictors. We will now look at "
    "each type with a data set from a different profession.")
types = [
    ("Simple", "One predictor\nStraight line", "y = b₀ + b₁x"),
    ("Multiple", "Two or more predictors\nFlat plane", "y = b₀ + b₁x₁ + b₂x₂ + …"),
    ("Polynomial", "Powers of a predictor\nCurve", "y = b₀ + b₁x + b₂x²"),
    ("Regularised", "Ridge / Lasso\nPenalised coefficients", "least squares + λ × penalty"),
]
for i, (name, desc, eq) in enumerate(types):
    left = Inches(0.6 + i * 3.1)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(2.9), Inches(3.9))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = BLUE
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(2.9), Inches(0.8))
    head.fill.solid(); head.fill.fore_color.rgb = BLUE; head.line.fill.background()
    add_text(s, left, Inches(1.6), Inches(2.9), Inches(0.8), name, size=20, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.15), Inches(2.6), Inches(2.6), Inches(1.5), desc.split("\n"),
             size=15, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(s, left + Inches(0.15), Inches(4.2), Inches(2.6), Inches(1.1), eq,
             size=14, color=NAVY, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, Inches(0.6), Inches(5.8), Inches(12), Inches(1),
         ["All are “linear” because the outcome is a linear combination of the coefficients.",
          "Categorical predictors (e.g. manual vs automatic) enter any of them as 0 / 1 dummy variables."],
         size=15, color=GRAY, align=PP_ALIGN.CENTER)

s = content_slide("Simple Linear Regression : Agriculture", "Part 2 : Types",
    "One predictor, one slope. The whole story fits on a single scatter plot, which makes it "
    "ideal for explaining results to farmers or managers who are not statisticians. Its "
    "weakness is that it ignores every other cause: rainfall, soil quality, seed variety. If "
    "those are correlated with fertilizer use, the slope will be biased. Use simple regression "
    "when there is a single dominant driver or as a first look before building a bigger model.")
add_bullets(s, Inches(0.6), Inches(1.4), Inches(6.3), Inches(5.4), [
    "**Data set**: crop yield (bushels / acre) vs fertilizer applied (kg / ha) on 30 field plots",
    "**Model**:  yield = b₀ + b₁ × fertilizer",
    "**Reading**: b₁ = extra bushels per extra kg of fertilizer",
    "**Strengths**",
    ("Easy to draw, easy to explain to non-statisticians", 1),
    ("Only two numbers to interpret", 1),
    "**Weaknesses**",
    ("Ignores rainfall, soil, seed variety", 1),
    ("Slope is biased if a missing cause is linked to fertilizer use", 1),
], size=17)
add_code(s, Inches(7.2), Inches(1.4), Inches(5.6), Inches(2.3),
         "farm  <- read.csv(\"farm_plots.csv\")\nmodel <- lm(yield ~ fertilizer, data = farm)\nsummary(model)\nplot(farm$fertilizer, farm$yield)\nabline(model, col = \"red\")")
add_bullets(s, Inches(7.2), Inches(3.9), Inches(5.6), Inches(2.8), [
    "Best for: a **single dominant driver**, quick first look, teaching",
    "In the summative Task 2 data, fertilizer type is categorical, so ANOVA is the matching tool; with a **continuous** fertilizer amount, simple regression is the tool",
], size=15)

s = content_slide("Multiple Linear Regression : Real Estate", "Part 2 : Types",
    "This is the workhorse of business analytics. The key phrase is holding the others "
    "constant: b1 is the price of one extra square foot for houses with the same number of "
    "bedrooms, the same age and the same distance. Use adjusted R-squared rather than R-squared "
    "to compare models, because R-squared always rises when you add a predictor even if it is "
    "useless. Watch for multicollinearity: area and bedrooms are strongly related, so the model "
    "struggles to separate their effects. Variable selection by AIC drops predictors that do not "
    "earn their place. In the car example from the summative task, stepwise selection reduced "
    "ten predictors to three: weight, quarter-mile time and transmission.")
add_bullets(s, Inches(0.6), Inches(1.4), Inches(6.3), Inches(5.4), [
    "**Data set**: house sale prices with area, bedrooms, age, distance to city centre",
    "**Model**: price = b₀ + b₁·area + b₂·bedrooms + b₃·age + b₄·distance",
    "Each slope = effect of that predictor **holding the others constant**",
    "Compare models with **adjusted R²** and **AIC**, not plain R²",
    "Danger: **multicollinearity** (area and bedrooms move together)",
    "Fix: variable selection, e.g. backward stepwise by AIC",
], size=17)
add_code(s, Inches(7.2), Inches(1.4), Inches(5.6), Inches(2.0),
         "model <- lm(price ~ area + bedrooms + age + distance,\n            data = houses)\nsummary(model)$adj.r.squared\nstep(model, direction = \"backward\")   # selection")
add_table(s, Inches(7.2), Inches(3.6), Inches(5.6), Inches(2.9), [
    ["Car data (summative Task 1)", "Adj. R²", "AIC"],
    ["mpg ~ wt", "0.745", "166.0"],
    ["mpg ~ all 10 predictors", "0.807", "163.7"],
    ["mpg ~ wt + qsec + am  (stepwise)", "0.834", "154.1"],
], col_widths=[3.4, 1.1, 1.1], size=13)

s = content_slide("Polynomial Regression : Pharmacology", "Part 2 : Types",
    "The dose-response curve is the classic case. Response rises with dose, then flattens or "
    "falls. A straight line cannot capture that, but adding a squared term can. A negative b2 "
    "gives an inverted U: more drug helps up to a point and then harms. The model is still "
    "linear in its coefficients, so lm handles it. Two warnings: always plot the fitted curve "
    "over the data, and never extrapolate a polynomial beyond the observed doses because the "
    "curve can shoot off in either direction.")
add_image(s, "polynomial.png", Inches(0.5), Inches(1.4), width=Inches(6.6))
add_bullets(s, Inches(7.4), Inches(1.4), Inches(5.5), Inches(2.6), [
    "**Data set**: drug dose (mg) vs patient response in a clinical trial",
    "Response rises, then flattens or falls: a **curve**, not a line",
    "**Model**: response = b₀ + b₁·dose + b₂·dose²",
    "Still linear in the coefficients, so **lm()** fits it",
], size=16)
add_code(s, Inches(7.4), Inches(4.1), Inches(5.5), Inches(1.5),
         "model <- lm(response ~ dose + I(dose^2), data = trial)\n# or   lm(response ~ poly(dose, 2), data = trial)")
add_bullets(s, Inches(7.4), Inches(5.7), Inches(5.5), Inches(1.2), [
    "Risk: high degrees **overfit** and behave wildly outside the data range",
], size=16)

s = content_slide("Compare and Contrast the Types", "Part 2 : Types",
    "Contrast on three axes: number of predictors, shape of the fitted relationship, and what "
    "can go wrong. Simple regression is the most interpretable but the least realistic. "
    "Multiple regression is realistic but needs care with correlated predictors. Polynomial "
    "handles curves but overfits easily. Ridge and Lasso trade a little bias for much lower "
    "variance when there are dozens of predictors, at the cost of coefficients that are harder "
    "to interpret. The rule: choose the simplest model that passes the assumption checks and "
    "has the lowest AIC or highest adjusted R-squared.")
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.4), [
    ["Type", "Predictors", "Shape", "Example data set", "Best for", "Main risk"],
    ["Simple", "1", "Straight line", "Fertilizer vs yield (agriculture)", "Explaining, quick checks", "Omitted variables"],
    ["Multiple", "2 or more", "Plane", "House prices (real estate)", "Controlling for several causes", "Multicollinearity, overfitting"],
    ["Polynomial", "1 + its powers", "Curve", "Dose vs response (pharmacology)", "Non-linear trends", "Overfitting, bad extrapolation"],
    ["Ridge / Lasso", "Many", "Plane, shrunk", "50 marketing channels (business)", "Many correlated predictors", "Harder to interpret"],
], col_widths=[1.5, 1.5, 1.5, 3.1, 2.5, 2.2], size=13)
add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9),
         "Rule of thumb: pick the simplest model that passes the LINE checks and has the lowest AIC / highest adjusted R².",
         size=17, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# =====================================================================
#  PART 3 : LINEAR vs LOGISTIC  (slides 10-13)
# =====================================================================
s = content_slide("Linear vs Logistic Regression : the Core Difference", "Part 3 : Linear vs Logistic",
    "The single question to ask is: is my outcome a number, or a yes-no? A number means linear "
    "regression. A yes-no means logistic regression. Why not draw a straight line through "
    "zeros and ones? Because the line predicts values below zero and above one, which are not "
    "probabilities. The logistic curve, an S shape, squeezes any value into the zero to one "
    "range. Logistic regression is fitted by maximum likelihood rather than least squares, and "
    "its coefficients are read as odds ratios after taking the exponential.")
add_image(s, "linear_vs_logistic.png", Inches(0.5), Inches(1.4), width=Inches(6.4))
add_table(s, Inches(7.1), Inches(1.4), Inches(5.8), Inches(5.3), [
    ["", "Linear", "Logistic"],
    ["Outcome y", "Numeric (price, mpg)", "Yes / No  (0 / 1)"],
    ["Predicts", "Value of y", "Probability y = 1"],
    ["Equation", "y = b₀ + b₁x", "log(p / (1 − p)) = b₀ + b₁x"],
    ["Shape", "Straight line", "S-curve between 0 and 1"],
    ["Fitting", "Least squares", "Maximum likelihood"],
    ["Coefficient", "Change in y per unit x", "exp(b) = odds ratio"],
    ["R function", "lm()", "glm(family = binomial)"],
    ["Evaluate", "R², RMSE", "Accuracy, confusion matrix"],
], col_widths=[1.4, 2.1, 2.3], size=12)

s = content_slide("Real-Life Example : a Bank and its Customers", "Part 3 : Linear vs Logistic",
    "Same bank, same customers, same four predictors. What decides the method is the type of "
    "outcome. For the spending question the model returns a rupee amount with a confidence "
    "interval. For the default question the model returns a probability, say 0.23. The bank "
    "then chooses a cut-off: above 0.20 reject, below approve. Moving the cut-off trades false "
    "approvals against false rejections. Logistic coefficients are read as odds ratios: exp of "
    "b equal to 1.5 means each extra unit of x multiplies the odds of default by 1.5.")
add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.6),
         "A bank records for each customer:  income,  age,  credit score,  existing debt", size=18, bold=True, color=NAVY)
for i, (title, q, out, code, color) in enumerate([
    ("Question A  →  LINEAR", "“How much will this customer spend on their card next month?”",
     "Outcome = amount in rupees (a number)", "lm(spend ~ income + age + score + debt, data = bank)", BLUE),
    ("Question B  →  LOGISTIC", "“Will this customer default on the loan?”",
     "Outcome = default Yes / No (a category)", "glm(default ~ income + age + score + debt,\n    data = bank, family = binomial)", RED),
]):
    left = Inches(0.6 + i * 6.2)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.2), Inches(5.9), Inches(3.9))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = color
    head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.2), Inches(5.9), Inches(0.7))
    head.fill.solid(); head.fill.fore_color.rgb = color; head.line.fill.background()
    add_text(s, left, Inches(2.2), Inches(5.9), Inches(0.7), title, size=18, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.2), Inches(3.05), Inches(5.5), Inches(1.4), [q, out], size=15, color=BLACK)
    add_code(s, left + Inches(0.2), Inches(4.5), Inches(5.5), Inches(1.1), code, size=12)
add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.6),
         "Same predictors, same data. The TYPE OF OUTCOME decides the method.",
         size=17, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

s = content_slide("Worked Example in R : Predicting Diabetes", "Part 3 : Linear vs Logistic",
    "The Pima Indians Diabetes data has 768 women and records whether each developed diabetes. "
    "Each extra unit of glucose multiplies the odds of diabetes by 1.03, so about twenty extra "
    "units roughly doubles the odds. A new patient with glucose 150, BMI 33 and age 45 gets a "
    "predicted probability of 0.63. On a held-out test set the model is 75 percent accurate. "
    "It is better at ruling out diabetes, specificity 87 percent, than at catching it, "
    "sensitivity 53 percent. Lowering the cut-off from 0.5 to 0.3 would catch more cases at "
    "the cost of more false alarms. That cut-off trade-off does not exist in linear regression, "
    "which is another practical difference. If instead we wanted to predict the glucose level "
    "itself, a number, we would go back to lm.")
add_image(s, "glucose_box.png", Inches(0.5), Inches(1.4), width=Inches(5.2))
add_code(s, Inches(5.9), Inches(1.4), Inches(7.0), Inches(2.6),
         "model <- glm(diabetes ~ glucose + bmi + age,\n             data = train, family = binomial)\nexp(coef(model))          # odds ratios\n#  glucose 1.03    bmi 1.09    age 1.03\n\npredict(model, type = \"response\",\n        newdata = data.frame(glucose = 150, bmi = 33, age = 45))\n#  0.63   ->  63 % probability of diabetes", size=12)
add_table(s, Inches(5.9), Inches(4.5), Inches(3.4), Inches(1.6), [
    ["Test-set result", "Value"],
    ["Accuracy", "75 %"],
    ["Sensitivity", "53 %"],
    ["Specificity", "87 %"],
], col_widths=[2.0, 1.4], size=13)
add_bullets(s, Inches(9.5), Inches(4.5), Inches(3.4), Inches(2.2), [
    "Cut-off 0.5 → lower it to catch more cases",
    "Predicting the glucose **level** instead would need **lm()**",
], size=14)

s = content_slide("Summary : Which Method, When?", "Part 3 : Summary",
    "Close by repeating the decision rule: look at your y first. A number means linear "
    "regression; choose the simplest type that fits and passes the LINE checks. A yes-no means "
    "logistic regression; read odds ratios and choose a cut-off. Then invite questions.")
add_bullets(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.2), [
    "**Linear regression**",
    ("Numeric outcome, straight-line relationship", 1),
    ("Read the slopes directly, check R² and residual plots", 1),
    "**Types**: simple, multiple, polynomial, regularised",
    ("Pick the simplest that passes the checks and has the lowest AIC", 1),
    "**Logistic regression**",
    ("Yes / No outcome, predicts a probability", 1),
    ("Read odds ratios, choose a cut-off, judge by a confusion matrix", 1),
], size=18)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.4))
box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.fill.background()
add_text(s, Inches(7.2), Inches(2.0), Inches(5.4), Inches(1.0), "Decision rule", size=22, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(7.2), Inches(2.9), Inches(5.4), Inches(3.2), [
    "Look at your  y  first.",
    "",
    "Is it a number?      →  lm()",
    "Is it Yes / No?       →  glm(family = binomial)",
    "",
    "Then check the assumptions before you trust a single p-value.",
], size=17, color=WHITE, align=PP_ALIGN.CENTER)

# =====================================================================
#  REFERENCES + THANK YOU  (not counted)
# =====================================================================
s = content_slide("References", "",
    "References in Harvard style. The mtcars and Pima data sets are the sources of all R "
    "output in this presentation.")
add_bullets(s, Inches(0.6), Inches(1.4), Inches(12.2), Inches(5.5), [
    "James, G., Witten, D., Hastie, T. and Tibshirani, R. (2021) An Introduction to Statistical Learning with Applications in R. 2nd edn. New York: Springer.",
    "Field, A., Miles, J. and Field, Z. (2012) Discovering Statistics Using R. London: Sage.",
    "Kutner, M. H., Nachtsheim, C. J., Neter, J. and Li, W. (2005) Applied Linear Statistical Models. 5th edn. New York: McGraw-Hill.",
    "Hosmer, D. W., Lemeshow, S. and Sturdivant, R. X. (2013) Applied Logistic Regression. 3rd edn. Hoboken: Wiley.",
    "R Core Team (2023) R: A Language and Environment for Statistical Computing. Vienna: R Foundation for Statistical Computing. Available at: https://www.R-project.org/",
    "Henderson, H. V. and Velleman, P. F. (1981) ‘Building multiple regression models interactively’, Biometrics, 37, pp. 391–411. (source of the mtcars data set)",
    "Smith, J. W. et al. (1988) ‘Using the ADAP learning algorithm to forecast the onset of diabetes mellitus’, Proceedings of the Symposium on Computer Applications and Medical Care, pp. 261–265. (source of the Pima Indians Diabetes data set)",
], size=15)

s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
add_text(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.5), "Thank you", size=54, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(1), "Questions and discussion", size=28,
         color=RGBColor(0xC9, 0xD3, 0xDF), align=PP_ALIGN.CENTER)
s.notes_slide.notes_text_frame.text = "Thank the audience and invite questions."

prs.save(OUT)
print("saved", OUT, "with", len(prs.slides), "slides")
